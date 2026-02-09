#!/usr/bin/env python3
"""Extract all text content from a PowerPoint file."""

import subprocess
import sys

# Ensure python-pptx is installed
try:
    from pptx import Presentation
except ImportError:
    subprocess.check_call([sys.executable, "-m", "pip", "install", "--break-system-packages", "python-pptx"])
    from pptx import Presentation

from pptx.enum.shapes import MSO_SHAPE_TYPE

def extract_table_text(table):
    """Extract text from a table shape with row/column formatting."""
    rows = []
    for row_idx, row in enumerate(table.rows):
        cells = []
        for cell in row.cells:
            cells.append(cell.text.strip())
        rows.append(cells)

    if not rows:
        return ""

    # Calculate column widths
    num_cols = max(len(r) for r in rows)
    col_widths = [0] * num_cols
    for row in rows:
        for i, cell in enumerate(row):
            for line in cell.split('\n'):
                col_widths[i] = max(col_widths[i], len(line))

    # Cap column widths at 80 chars
    col_widths = [min(w, 80) for w in col_widths]

    result = []
    separator = "+" + "+".join("-" * (w + 2) for w in col_widths) + "+"

    for row_idx, row in enumerate(rows):
        # Pad row if needed
        while len(row) < num_cols:
            row.append("")

        # Handle multi-line cells
        cell_lines = [cell.split('\n') for cell in row]
        max_lines = max(len(lines) for lines in cell_lines)

        if row_idx == 0:
            result.append(separator)

        for line_idx in range(max_lines):
            line_parts = []
            for col_idx, lines in enumerate(cell_lines):
                text = lines[line_idx] if line_idx < len(lines) else ""
                if len(text) > col_widths[col_idx]:
                    text = text[:col_widths[col_idx]-3] + "..."
                line_parts.append(text.ljust(col_widths[col_idx]))
            result.append("| " + " | ".join(line_parts) + " |")

        result.append(separator)

    return "\n".join(result)


def extract_group_text(shape):
    """Recursively extract text from grouped shapes."""
    texts = []
    if hasattr(shape, 'shapes'):
        for s in shape.shapes:
            if s.has_text_frame:
                for para in s.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        texts.append(text)
            if s.has_table:
                texts.append(extract_table_text(s.table))
            if hasattr(s, 'shapes'):
                texts.append(extract_group_text(s))
    return "\n".join(t for t in texts if t)


def main():
    filepath = "/Users/Paco/Documents/M2 GIP/INGENIE FISCALE DES PERSONNES PHYSIQUES & PATRIMOINE/Partie 4.pptx"

    print(f"Opening: {filepath}")
    prs = Presentation(filepath)
    print(f"Total slides: {len(prs.slides)}")
    print("=" * 80)

    for slide_num, slide in enumerate(prs.slides, 1):
        print(f"\n{'=' * 80}")
        print(f"SLIDE {slide_num}")
        print(f"{'=' * 80}")

        # Get slide layout name if available
        try:
            layout_name = slide.slide_layout.name
            print(f"[Layout: {layout_name}]")
        except:
            pass

        for shape_idx, shape in enumerate(slide.shapes):
            shape_name = shape.name if hasattr(shape, 'name') else f"Shape {shape_idx}"

            # Text frames (titles, text boxes, etc.)
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        print(text)

            # Tables
            if shape.has_table:
                print(f"\n[Table: {shape_name}]")
                print(extract_table_text(shape.table))
                print()

            # Grouped shapes
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                group_text = extract_group_text(shape)
                if group_text:
                    print(group_text)

        # Check for notes
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                print(f"\n[Speaker Notes]:")
                print(notes_text)

    print(f"\n{'=' * 80}")
    print("EXTRACTION COMPLETE")
    print(f"{'=' * 80}")


if __name__ == "__main__":
    main()
